from pptx import Presentation
from pptx.util import Inches, Pt
import os


def add_bullet_slide(pr, title, bullets):
    layout = pr.slide_layouts[1]
    slide = pr.slides.add_slide(layout)
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = b
        else:
            p = tf.add_paragraph()
            p.text = b
        p.level = 0


def main():
    pr = Presentation()

    # Slide 1 - Title
    title_slide_layout = pr.slide_layouts[0]
    slide = pr.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Gestion de Cimetière - Présentation"
    slide.placeholders[1].text = "Présentation du projet"

    # Slide sections
    add_bullet_slide(pr, "Objectif", [
        "Fournir une application de gestion des caveaux, réservations et concessions",
        "Gérer les utilisateurs, MFA et rôles (ADMIN, AGENT, SECRETAIRE, CLIENT)",
    ])

    add_bullet_slide(pr, "Introduction", [
        "Application full-stack: backend Django + API (django-ninja)",
        "Frontend desktop/web en Flet pour l'interface utilisateur",
    ])

    add_bullet_slide(pr, "Architecture", [
        "Backend: Django, Django Ninja pour l'API (cimetiere_api)",
        "Base de données: SQLite (dev) ou PostgreSQL (prod)",
        "Frontend: Flet (web) - single executable UI",
    ])

    add_bullet_slide(pr, "Fonctionnement", [
        "Inscription, login, MFA par email (code 6 chiffres)",
        "Gestion des caveaux, réservations, concessions, exhumations, transactions",
        "Tableau de bord avec stats et cartographie interactive",
    ])

    add_bullet_slide(pr, "Sécurité utilisée", [
        "Mots de passe hachés via django.contrib.auth.hashers",
        "MFA par code envoyé par email (TTL 5 minutes)",
        "Permissions basées sur le rôle (ADMIN/AGENT/SECRETAIRE/CLIENT)",
    ])

    # Modules installés (requirements)
    req_path = os.path.join('backend', 'requirements.txt')
    if not os.path.exists(req_path):
        req_path = 'requirements.txt'
    modules = []
    if os.path.exists(req_path):
        with open(req_path, 'r', encoding='utf-8') as f:
            for l in f:
                s = l.strip()
                if s and not s.startswith('#'):
                    modules.append(s)

    add_bullet_slide(pr, "Modules / Bibliothèques", modules[:20] or ["Voir requirements.txt"])

    add_bullet_slide(pr, "Rôles", [
        "ADMIN: accès complet, revenus et gestion utilisateurs",
        "AGENT: accès terrain, gestion caveaux",
        "SECRETAIRE: gestion réservations et clients",
        "CLIENT: consulter son profil et faire réservations",
    ])

    add_bullet_slide(pr, "Difficultés rencontrées", [
        "Migration de schéma entre versions et compatibilité SQLite/Postgres",
        "Configuration d'envoi d'emails en dev sans divulguer secrets",
        "Intégration Flet Web + appels API locaux",
    ])

    add_bullet_slide(pr, "Accès & Langages", [
        "Langages: Python 3.11+, Django 6, Flet pour UI",
        "Accès: interface web locale (Flet) et /admin/ Django pour super-admin",
    ])

    add_bullet_slide(pr, "Conclusion", [
        "Solution opérationnelle pour la gestion funéraire locale",
        "Extensible vers déploiement production (Postgres, WSGI/ASGI)",
    ])

    out_path = os.path.join('backend', 'Presentation_Gestion_Cimetiere.pptx')
    # ensure output directory exists
    out_dir = os.path.dirname(out_path)
    if out_dir and not os.path.exists(out_dir):
        os.makedirs(out_dir, exist_ok=True)
    pr.save(out_path)
    print(f"Présentation créée: {out_path}")


if __name__ == '__main__':
    main()
